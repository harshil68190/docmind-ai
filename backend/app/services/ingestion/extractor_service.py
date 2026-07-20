"""
Text extraction service.

Returns a list of `ExtractedPage` records rather than one flat string,
because citation-aware answers (Milestone 4) need to know which page a
chunk of text came from. PDF and PPTX have a natural page/slide concept;
DOCX and TXT don't -- python-docx exposes paragraphs, not rendered page
breaks -- so those formats return a single page with `page_number=None`.
Callers already have to handle a missing page number: the citation
requirement is "include the source file and page number whenever
available", which anticipates exactly this case.
"""
from dataclasses import dataclass
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

from app.core.exceptions import ExtractionFailedException


@dataclass(frozen=True)
class ExtractedPage:
    """One page's text (or, for page-less formats, the whole document's)."""

    page_number: int | None
    text: str


_EXTRACTORS = {}  # populated at the bottom of this file


class ExtractionService:
    def extract(self, file_path: str, file_extension: str) -> list[ExtractedPage]:
        extractor = _EXTRACTORS.get(file_extension)
        if extractor is None:
            raise ExtractionFailedException(
                f"No text extractor is registered for '{file_extension}' files."
            )
        try:
            raw_pages = extractor(file_path)
        except Exception as exc:  # noqa: BLE001 - any parser failure becomes a domain error
            raise ExtractionFailedException(
                f"Failed to extract text from this {file_extension} file: {exc}"
            ) from exc

        pages = [
            ExtractedPage(page.page_number, page.text.strip())
            for page in raw_pages
            if page.text.strip()
        ]
        if not pages:
            raise ExtractionFailedException(
                "No extractable text was found in this document "
                "(it may be a scanned/image-only file)."
            )
        return pages

    @staticmethod
    def _extract_pdf(file_path: str) -> list[ExtractedPage]:
        with fitz.open(file_path) as pdf:
            return [
                ExtractedPage(page_number=index + 1, text=page.get_text())
                for index, page in enumerate(pdf)
            ]

    @staticmethod
    def _extract_docx(file_path: str) -> list[ExtractedPage]:
        document = DocxDocument(file_path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        return [ExtractedPage(page_number=None, text=text)]

    @staticmethod
    def _extract_pptx(file_path: str) -> list[ExtractedPage]:
        presentation = Presentation(file_path)
        pages: list[ExtractedPage] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text_parts = [
                shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
            ]
            pages.append(
                ExtractedPage(page_number=slide_number, text="\n".join(text_parts))
            )
        return pages

    @staticmethod
    def _extract_txt(file_path: str) -> list[ExtractedPage]:
        text = Path(file_path).read_text(encoding="utf-8")
        return [ExtractedPage(page_number=None, text=text)]


_EXTRACTORS[".pdf"] = ExtractionService._extract_pdf
_EXTRACTORS[".docx"] = ExtractionService._extract_docx
_EXTRACTORS[".pptx"] = ExtractionService._extract_pptx
_EXTRACTORS[".txt"] = ExtractionService._extract_txt
